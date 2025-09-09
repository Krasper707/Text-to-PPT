
# app.py (reloaded)

import streamlit as st
import json
import requests
from pptx import Presentation
import io
from typing import Union # <-- IMPORT THIS FOR OLDER PYTHON VERSIONS
from pptx.enum.text import MSO_AUTO_SIZE
import base64

# --- PAGE CONFIGURATION ---
st.set_page_config(page_title="Your Text, Your Style", page_icon="✨", layout="centered")

# --- SESSION STATE ---
if 'models' not in st.session_state: st.session_state.models = []
if 'selected_model' not in st.session_state: st.session_state.selected_model = None

# =====================================================================================
# LOGIC FUNCTIONS (WITH COMPATIBLE TYPE HINTS)
# =====================================================================================

@st.cache_data(show_spinner="Fetching available models...")
def get_available_models(aipipe_token: str) -> list[str]:
    # This function remains the same.
    if not aipipe_token: return []
    url = "https://aipipe.org/openrouter/v1/models"
    headers = {"Authorization": f"Bearer {aipipe_token}"}
    try:
        response = requests.get(url, headers=headers)
        if response.status_code == 200:
            return sorted([model['id'] for model in response.json().get('data', [])])
        else:
            st.error(f"Failed to fetch models. Status: {response.status_code} - {response.text}")
            return []
    except Exception as e:
        st.error(f"An error occurred while fetching models: {e}")
        return []

def generate_slide_content(text: str, guidance: str, aipipe_token: str, model_name: str) -> Union[dict, None]: # <-- FIX
    """Uses an LLM to structure text and suggest visuals."""
    st.info(f"Sending request to model: {model_name}...")
    url = "https://aipipe.org/openrouter/v1/chat/completions"
    headers = {"Authorization": f"Bearer {aipipe_token}", "Content-Type": "application/json"}
    
    # system_prompt = """
    # You are an expert presentation creator. Your task is to analyze the following text and user guidance to structure it into a series of presentation slides.
    # - Break down the content logically. The number of slides should be reasonable for the text provided.
    # - For each slide, provide a `title`, `content` as bullet points, and `speaker_notes`.
    # - For each slide, also provide a `visual_suggestion`. This should be a brief, one-sentence description of an ideal image.
    # - If the slide content is purely abstract or does not need an image, the value for `visual_suggestion` MUST be the exact string "none".
    # - You MUST output your response ONLY as a single, valid JSON object following this exact schema:
    # {"slides": [{"title": "Slide Title", "content": ["Bullet point 1."], "speaker_notes": "Notes.", "visual_suggestion": "A picture of a computer."}]}
    # Do not include any other text, explanations, or markdown.
    # """
    system_prompt = """
    You are an expert presentation creator. Your task is to analyze the following text and user guidance to structure it into a series of presentation slides.
    - Break down the content logically.
    - For each slide, provide a `title`, `content`, `speaker_notes`, and a `visual_suggestion`.
    - For the `content` field, create a list of strings. To create nested bullet points, start a string with '- ' for a sub-bullet.
    - For `visual_suggestion`, describe an ideal image. If no image is needed, use the exact string "none".
    - You MUST output your response ONLY as a single, valid JSON object following this exact schema:
    {"slides": [{"title": "Main Point", "content": ["First bullet.", "- Sub-bullet.", "Second bullet."], "speaker_notes": "Notes.", "visual_suggestion": "A relevant image."}]}
    """

    user_content = f"### USER GUIDANCE:\n{guidance}\n\n### SOURCE TEXT:\n{text}"
    payload = { "model": model_name, "messages": [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_content}], "response_format": {"type": "json_object"} }
    
    try:
        response = requests.post(url, headers=headers, json=payload, timeout=180)
        if response.status_code == 200:
            return json.loads(response.json()['choices'][0]['message']['content'])
        else:
            st.error(f"API Error: Status {response.status_code} - Response: {response.text}")
            return None
    except Exception as e:
        st.error(f"An exception occurred making the API request: {e}")
        return None

def analyze_template(template_file) -> tuple[Union[Presentation, None], dict, list]:
    """
    Analyzes a template to find layouts and extract a *clean* image bank.
    Only extracts images from actual Picture placeholders to avoid background art.
    """
    try:
        prs = Presentation(template_file)
        layouts, image_bank = {}, []
        
        # --- MORE ROBUST LAYOUT CLASSIFICATION ---
        for layout in prs.slide_layouts:
            placeholders = {p.placeholder_format.type for p in layout.placeholders}
            has_title = any(p.placeholder_format.type == 1 for p in layout.placeholders)
            has_body = any(p.placeholder_format.type == 2 for p in layout.placeholders)
            has_picture = any(p.placeholder_format.type == 18 for p in layout.placeholders)
            
            if has_title and has_body and has_picture: layouts['content_with_image'] = layout
            elif has_title and has_body: layouts['content_only'] = layout

        layouts['title'] = prs.slide_layouts[0]
        if 'content_only' not in layouts: layouts['content_only'] = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        if 'content_with_image' not in layouts: layouts['content_with_image'] = layouts.get('content_only')

        # --- CRITICAL FIX: Only extract images from Picture Placeholders ---
        for slide in prs.slides:
            for shape in slide.shapes:
                # We check if the shape is a placeholder AND its type is PICTURE (18)
                if shape.is_placeholder and shape.placeholder_format.type == 18 and hasattr(shape, 'image'):
                    image_bank.append(io.BytesIO(shape.image.blob))
        
        if image_bank:
            st.success(f"Template analyzed: Found {len(image_bank)} reusable content images.")
        else:
            st.info("Template analyzed: No reusable content images found. Will create a text-only presentation.")

        return prs, layouts, image_bank
    except Exception as e:
        st.error(f"Could not read the PowerPoint file. It may be corrupt. Error: {e}")
        return None, {}, []


def create_presentation(slide_data: dict, template_file) -> Union[io.BytesIO, None]:
    """
    Builds the presentation with a definitive, multi-step hierarchical method
    for finding the correct content placeholder.
    """
    prs, layouts, image_bank = analyze_template(template_file)
    if not prs: return None

    # Delete original slides to ensure a clean slate
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    slides_list = slide_data.get('slides', [])
    image_idx = 0
    
    try:
        if not slides_list:
             st.error("AI returned an empty list of slides.")
             return None

        # --- CLEAN TITLE SLIDE CREATION ---
        slide_info = slides_list[0]
        slide = prs.slides.add_slide(layouts['title'])
        if slide.shapes.title:
            slide.shapes.title.text = slide_info.get('title', 'Presentation Title')

        # --- DEFINITIVE CONTENT SLIDE CREATION ---
        for slide_info in slides_list[1:]:
            visual_suggestion = slide_info.get('visual_suggestion', 'none').lower()
            use_image_layout = (visual_suggestion != 'none' and image_bank)
            layout = layouts['content_with_image'] if use_image_layout else layouts['content_only']
            slide = prs.slides.add_slide(layout)
            
            if slide.shapes.title:
                slide.shapes.title.text = slide_info.get('title', '')
            
            # --- HIERARCHICAL BODY SHAPE FINDING (v3 - Definitive) ---
            body_shape = None
            
            # 1. Ideal Case: Find the official "Body" placeholder (type 2)
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:
                    body_shape = shape
                    break
            
            # 2. Good Case: If no official body, find a generic content placeholder by elimination.
            if not body_shape:
                for shape in slide.placeholders:
                    # Exclude Title (1), Footer (14), Date (15), Slide Number (13)
                    if shape.placeholder_format.type not in [1, 14, 15, 13]:
                        if hasattr(shape, 'text_frame'): # Make sure it can hold text
                            body_shape = shape
                            break
            
            # 3. Last Resort Fallback: If still nothing, find the largest text box that isn't the title.
            if not body_shape:
                potential_shapes = [s for s in slide.placeholders if s.has_text_frame and s != slide.shapes.title]
                if potential_shapes:
                    body_shape = max(potential_shapes, key=lambda s: s.width * s.height)

            # --- POPULATE THE FOUND BODY SHAPE ---
            if body_shape:
                tf = body_shape.text_frame
                tf.clear()
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                tf.word_wrap = True
                content = slide_info.get('content', [])
                if isinstance(content, dict): content = list(content.values())
                if isinstance(content, list):
                    for point in content:
                        p = tf.add_paragraph()
                        if point.strip().startswith('- '):
                            p.text = point.strip().lstrip('- ').strip()
                            p.level = 1
                        else:
                            p.text = point
                            p.level = 0
            else:
                st.warning(f"Could not find a suitable content placeholder on slide titled '{slide_info.get('title', '')}'. This slide will be empty.")
            
            # --- IMAGE AND NOTES LOGIC (remains the same) ---
            if use_image_layout:
                pic_placeholder = next((s for s in slide.placeholders if s.placeholder_format.type == 18), None)
                if pic_placeholder:
                    img_stream = image_bank[image_idx % len(image_bank)]
                    img_stream.seek(0)
                    pic_placeholder.insert_picture(img_stream)
                    image_idx += 1
            
            if 'speaker_notes' in slide_info:
                slide.notes_slide.notes_text_frame.text = slide_info['speaker_notes']

        powerpoint_stream = io.BytesIO()
        prs.save(powerpoint_stream)
        powerpoint_stream.seek(0)
        return powerpoint_stream
    except Exception as e:
        st.error(f"An error occurred while building slides: {e}")
        return None

def create_revealjs_presentation(slide_data:dict,image_bank:list):
    """Creates a simple RevealJS HTML presentation."""
    try:
        slides_html = ""
        image_idx = 0
        slides_list = slide_data.get('slides', [])

        if not slides_list:
            st.error("Cannot create Reveal.js slides: AI data is empty.")
            return None

        # Create the HTML for each slide as a <section>
        for slide_info in slides_list:
            title = slide_info.get('title', '')
            content = slide_info.get('content', [])
            notes = slide_info.get('speaker_notes', '')
            visual_suggestion = slide_info.get('visual_suggestion', 'none').lower()

            if isinstance(content, dict): content = list(content.values())

            # Build the HTML for the bullet points
            content_html = "<ul>\n"
            for point in content:
                if point.strip().startswith('- '):
                    content_html += f"<ul><li>{point.strip().lstrip('- ').strip()}</li></ul>\n"
                else:
                    content_html += f"<li>{point}</li>\n"
            content_html += "</ul>"
            
            # Embed an image if suggested and available
            image_html = ""
            if visual_suggestion != 'none' and image_bank:
                img_stream = image_bank[image_idx % len(image_bank)]
                img_stream.seek(0)
                b64_img = base64.b64encode(img_stream.read()).decode('utf-8')
                image_html = f'<img class="r-stretch" src="data:image/png;base64,{b64_img}" style="max-height: 450px; margin: auto;">'
                image_idx += 1
            
            notes_html = f'<aside class="notes">{notes}</aside>' if notes else ''

            # Assemble the final HTML for this slide
            slides_html += f"""
            <section>
                <h2>{title}</h2>
                {content_html}
                {image_html}
                {notes_html}
            </section>
            """

        # Wrap the slides in the full Reveal.js HTML boilerplate
        full_html = f"""
        <!doctype html>
        <html>
            <head>
                <meta charset="utf-8">
                <title>{slides_list[0].get('title', 'Presentation')}</title>
                <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/reset.min.css">
                <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/reveal.min.css">
                <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/theme/black.min.css" id="theme">
            </head>
            <body>
                <div class="reveal">
                    <div class="slides">
                        {slides_html}
                    </div>
                </div>
                <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/reveal.min.js"></script>
                <script>
                    Reveal.initialize({{ hash: true, plugins: [] }});
                </script>
            </body>
        </html>
        """
        return full_html

    except Exception as e:
        st.error(f"An error occurred while building the Reveal.js HTML: {e}")
        return None

# =====================================================================================
# UI COMPONENTS & WORKFLOW
# =====================================================================================
st.title("✨ Your Text, Your Style")
st.subheader("Auto-Generate a Presentation from Text")

with st.expander("Step 1: Provide Your Source Text", expanded=True):
    source_text = st.text_area("Paste text here.", height=300, label_visibility="collapsed")

# with st.expander("Step 2: Configure the AI", expanded=True):
#     guidance = st.text_input("Optional guidance", placeholder="e.g., 'Turn into a 5-slide investor pitch deck'")
#     aipipe_token = st.text_input("Enter your AI Pipe Token", type="password")
    
#     if st.button("Load Available Models"):
#         st.session_state.models = get_available_models(aipipe_token)
#         if not st.session_state.models:
#             st.warning("Could not load models. Check your token.")

#     if st.session_state.models:
#         st.session_state.selected_model = st.selectbox(
#             "Choose a model:",
#             options=st.session_state.models,
#             index=st.session_state.models.index("anthropic/claude-3-5-sonnet-20240620") if "anthropic/claude-3-5-sonnet-20240620" in st.session_state.models else 0
#         )

with st.expander("Step 2: Configure the AI", expanded=True):
    # --- THIS IS THE NEW FEATURE ---
    guidance_templates = {
        "Default (Flexible)": "Break this down into a logical series of presentation slides.",
        "Investor Pitch Deck": "Structure this as a compelling 10-slide investor pitch deck, focusing on Problem, Solution, Market Size, Product, Business Model, and Team.",
        "Technical Deep Dive": "Convert this technical document into a presentation for an engineering audience. Focus on architecture, data flow, and key algorithms.",
        "Marketing Plan": "Turn this into a marketing plan presentation. Create slides for Target Audience, Key Messaging, Channels, and KPIs.",
        "Project Update": "Summarize this into a project update presentation with slides for 'What We Accomplished', 'Challenges Faced', and 'Next Steps'."
    }
    
    if 'guidance_text' not in st.session_state:
        st.session_state.guidance_text = guidance_templates["Default (Flexible)"]

    def update_guidance():
        st.session_state.guidance_text = guidance_templates[st.session_state.guidance_choice]

    st.selectbox(
        "Choose a guidance template (for tone and style):",
        options=guidance_templates.keys(),
        key="guidance_choice",
        on_change=update_guidance
    )
    
    guidance = st.text_input(
        "Or customize your own guidance:",
        key="guidance_text"
    )
    # --- END OF NEW FEATURE ---
    
    aipipe_token = st.text_input("Enter your AI Pipe Token", type="password")
    
    if st.button("Load Available Models"):
        st.session_state.models = get_available_models(aipipe_token)
        if not st.session_state.models:
            st.warning("Could not load models. Check your token.")

    if st.session_state.models:
        st.session_state.selected_model = st.selectbox(
            "Choose a model:",
            options=st.session_state.models,
            index=st.session_state.models.index("anthropic/claude-3-5-sonnet-20240620") if "anthropic/claude-3-5-sonnet-20240620" in st.session_state.models else 0
        )

with st.expander("Step 3: Upload Your Template", expanded=True):
    uploaded_template = st.file_uploader("Upload .pptx or .potx", type=['pptx', 'potx'], label_visibility="collapsed")

st.divider()

choice=st.radio("Which output would you prefer?",options=["PPTX file","RevealJS slides"],horizontal=True)
# if st.button("🚀 Generate Presentation", type="primary", use_container_width=True):
#     if not source_text or not aipipe_token or not st.session_state.selected_model or not uploaded_template:
#         st.warning("Please complete all steps before generating.")
#     else:
#         with st.spinner("🤖 AI is structuring content and suggesting visuals..."):
#             slide_data = generate_slide_content(source_text, guidance, aipipe_token, st.session_state.selected_model)
#         if slide_data:
#             st.success("✅ AI content generated!")
#             with st.spinner("🎨 Applying styles and reusing images..."):
#                 powerpoint_file = create_presentation(slide_data, uploaded_template)
#             if powerpoint_file:
#                 st.success("🎉 Your presentation is ready!")
#                 st.download_button(
#                     label="📥 Download Presentation (.pptx)",
#                     data=powerpoint_file,
#                     file_name="generated_presentation.pptx",
#                     mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
#                     use_container_width=True
#                 )
if st.button("🚀 Generate Presentation", type="primary", use_container_width=True):
    # --- UPGRADED VALIDATION ---
    # Basic checks that are always required
    if not source_text or not aipipe_token or not st.session_state.selected_model:
        st.warning("Please complete Steps 1 and 2 before generating.")
    # PowerPoint-specific check
    elif choice == "PPTX file" and not uploaded_template:
        st.warning("Please upload a PowerPoint template in Step 3 for PPTX generation.")
    else:
        with st.spinner("🤖 AI is structuring content and suggesting visuals..."):
            slide_data = generate_slide_content(source_text, guidance, aipipe_token, st.session_state.selected_model)
        
        if slide_data:
                    st.success("✅ AI content generated!")

                    # --- UPDATED LOGIC with Safety Nets ---

                    if choice == "PPTX file":
                        with st.spinner("🎨 Applying styles and building PowerPoint..."):
                            output_file = create_presentation(slide_data, uploaded_template)
                        
                        if output_file:
                            st.success("🎉 Your PowerPoint presentation is ready!")
                            st.download_button(
                                label="📥 Download Presentation (.pptx)",
                                data=output_file,
                                file_name="generated_presentation.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True
                            )
                        # --- THIS IS THE NEW SAFETY NET ---
                        else:
                            st.error("Failed to build the PowerPoint file. Please check for any errors above, the template file might be incompatible.")

                    elif choice == "RevealJS slides":
                        with st.spinner("🌐 Weaving HTML and embedding images for Reveal.js..."):
                            image_bank = []
                            if uploaded_template:
                                _, _, image_bank = analyze_template(uploaded_template)
                            output_file = create_revealjs_presentation(slide_data, image_bank)

                        if output_file:
                            st.success("🎉 Your Reveal.js slideshow is ready!")
                            st.download_button(
                                label="📥 Download Slideshow (.html)",
                                data=output_file,
                                file_name="generated_slideshow.html",
                                mime="text/html",
                                use_container_width=True
                            )
                        # --- THIS IS THE NEW SAFETY NET ---
                        else:
                            st.error("Failed to build the Reveal.js slideshow. Please check for any errors reported above.")
